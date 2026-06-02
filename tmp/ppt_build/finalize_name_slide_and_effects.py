from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
import shutil

from lxml import etree
from pptx import Presentation
from pptx.util import Pt


PPTX = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")
TMP_DIR = Path(r"E:\Github\Roguelike_Project\tmp\ppt_build\transition_tmp")
FONT = "Times New Roman"
BODY_SIZE = Pt(20)
SUBTITLE_SIZE = Pt(24)


def set_paragraph(paragraph, text, size=BODY_SIZE, bold=False):
    paragraph.clear()
    paragraph.line_spacing = 1.05
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(10)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold


def polish_text():
    prs = Presentation(str(PPTX))

    # Slide 15 after insertion: the map slide.
    slide = prs.slides[14]
    set_paragraph(slide.shapes[1].text_frame.paragraphs[0], "B\u1ea3n \u0111\u1ed3", size=SUBTITLE_SIZE, bold=True)
    set_paragraph(
        slide.shapes[1].text_frame.paragraphs[1],
        "- B\u1ea3n \u0111\u1ed3 thay \u0111\u1ed5i theo t\u1eebng giai \u0111o\u1ea1n \u0111\u1ec3 t\u1ea1o c\u1ea3m gi\u00e1c m\u1edbi trong qu\u00e1 tr\u00ecnh ch\u01a1i.",
    )
    set_paragraph(
        slide.shapes[1].text_frame.paragraphs[2],
        "- M\u1ed7i khi ho\u00e0n th\u00e0nh \u0111\u1ee3t c\u00f3 boss, tr\u00f2 ch\u01a1i chuy\u1ec3n sang b\u1ea3n \u0111\u1ed3 m\u1edbi.",
    )

    prs.save(str(PPTX))


def ensure_fade_transitions():
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    if TMP_DIR.exists():
        shutil.rmtree(TMP_DIR)
    TMP_DIR.mkdir(parents=True, exist_ok=True)
    tmp_pptx = TMP_DIR / PPTX.name

    with ZipFile(PPTX, "r") as zin, ZipFile(tmp_pptx, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                root = etree.fromstring(data)
                trans = root.find(f"{{{ns}}}transition")
                if trans is None:
                    trans = etree.Element(f"{{{ns}}}transition")
                    clr_map = root.find(f"{{{ns}}}clrMapOvr")
                    if clr_map is not None:
                        root.insert(root.index(clr_map) + 1, trans)
                    else:
                        root.insert(0, trans)
                trans.attrib.pop("advTm", None)
                trans.attrib.pop("advClick", None)
                for child in list(trans):
                    trans.remove(child)
                trans.append(etree.Element(f"{{{ns}}}fade"))
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(item, data)

    shutil.move(str(tmp_pptx), str(PPTX))
    shutil.rmtree(TMP_DIR, ignore_errors=True)


def main():
    polish_text()
    ensure_fade_transitions()
    print(PPTX)


if __name__ == "__main__":
    main()
