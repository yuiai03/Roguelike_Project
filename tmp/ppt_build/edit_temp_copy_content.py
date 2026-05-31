from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


PPTX_PATH = Path(r"E:\Github\Roguelike_Project\output\ppt\temp - Copy.pptx")
FONT_NAME = "Times New Roman"
BODY_SIZE = Pt(20)
PARAGRAPH_SPACE_AFTER = Pt(10)


def replace_text(shape, text, *, size=BODY_SIZE, bold=False, align=PP_ALIGN.LEFT):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = 1.05
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = size
    run.font.bold = bold


def set_bullet_frame(shape, lines):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.05
        paragraph.space_before = Pt(0)
        paragraph.space_after = PARAGRAPH_SPACE_AFTER
        run = paragraph.add_run()
        run.text = f"- {line}"
        run.font.name = FONT_NAME
        run.font.size = BODY_SIZE


def _capture_run_style(paragraph):
    run = next((item for item in paragraph.runs if item.text), None)
    if run is None:
        return {
            "font_name": FONT_NAME,
            "font_size": BODY_SIZE,
            "bold": False,
            "italic": False,
            "underline": False,
            "color": None,
        }

    color = None
    try:
        if run.font.color is not None and run.font.color.rgb is not None:
            color = run.font.color.rgb
    except AttributeError:
        color = None

    return {
        "font_name": run.font.name or FONT_NAME,
        "font_size": run.font.size or BODY_SIZE,
        "bold": run.font.bold,
        "italic": run.font.italic,
        "underline": run.font.underline,
        "color": color,
    }


def normalize_multiline_shape(shape):
    if not getattr(shape, "has_text_frame", False):
        return False

    paragraphs = shape.text_frame.paragraphs
    if len(paragraphs) <= 1:
        return False

    entries = []
    had_empty_paragraph = False
    for paragraph in paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            had_empty_paragraph = True
            continue

        entries.append(
            {
                "text": text,
                "alignment": paragraph.alignment,
                "style": _capture_run_style(paragraph),
            }
        )

    if not entries:
        return False

    needs_normalization = had_empty_paragraph or any(
        paragraph.space_after != PARAGRAPH_SPACE_AFTER or paragraph.space_before != Pt(0) or paragraph.line_spacing != 1.05
        for paragraph in paragraphs
        if "".join(run.text for run in paragraph.runs).strip()
    )
    if not needs_normalization:
        return False

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, entry in enumerate(entries):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = entry["alignment"] or PP_ALIGN.LEFT
        paragraph.line_spacing = 1.05
        paragraph.space_before = Pt(0)
        paragraph.space_after = PARAGRAPH_SPACE_AFTER

        run = paragraph.add_run()
        run.text = entry["text"]
        run.font.name = entry["style"]["font_name"]
        run.font.size = entry["style"]["font_size"]
        run.font.bold = entry["style"]["bold"]
        run.font.italic = entry["style"]["italic"]
        run.font.underline = entry["style"]["underline"]
        if entry["style"]["color"] is not None:
            run.font.color.rgb = entry["style"]["color"]

    return True


def main():
    prs = Presentation(str(PPTX_PATH))

    slide6 = prs.slides[5]
    replace_text(
        slide6.shapes[2],
        "D\u1ef1ng scene, gameplay v\u00e0 giao di\u1ec7n.",
    )
    replace_text(
        slide6.shapes[12],
        "M\u00f4i tr\u01b0\u1eddng l\u1eadp tr\u00ecnh v\u00e0 ch\u1ec9nh s\u1eeda m\u00e3 ngu\u1ed3n.",
    )

    slide11 = prs.slides[10]
    replace_text(
        slide11.shapes[1],
        "Kinh nghi\u1ec7m, l\u00ean c\u1ea5p v\u00e0 buff",
        size=BODY_SIZE,
        bold=True,
    )
    set_bullet_frame(
        slide11.shapes[3],
        [
            "Ng\u01b0\u1eddi ch\u01a1i nh\u1eadn kinh nghi\u1ec7m khi ti\u00eau di\u1ec7t enemy.",
            "Khi \u0111\u1ee7 kinh nghi\u1ec7m, nh\u00e2n v\u1eadt s\u1ebd l\u00ean c\u1ea5p.",
            "M\u1ed7i l\u1ea7n l\u00ean c\u1ea5p, ng\u01b0\u1eddi ch\u01a1i \u0111\u01b0\u1ee3c ch\u1ecdn buff \u0111\u1ec3 t\u0103ng s\u1ee9c m\u1ea1nh.",
            "C\u00e1c buff gi\u00fap thay \u0111\u1ed5i c\u00e1ch ch\u01a1i v\u00e0 h\u1ed7 tr\u1ee3 v\u01b0\u1ee3t qua c\u00e1c wave kh\u00f3 h\u01a1n.",
        ],
    )

    slide13 = prs.slides[12]
    replace_text(
        slide13.shapes[1],
        "Challenge v\u00e0 b\u1eaft \u0111\u1ea7u tr\u1eadn",
        size=BODY_SIZE,
        bold=True,
    )

    for slide_index in range(2, 16):
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            normalize_multiline_shape(shape)

    prs.save(str(PPTX_PATH))
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
