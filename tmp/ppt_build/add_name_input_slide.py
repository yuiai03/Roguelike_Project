from copy import deepcopy
from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
import shutil
import tempfile

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


PPTX = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")
IMAGE = Path(r"E:\Github\Roguelike_Project\output\ppt\assets_thaivanhao\doc_image_03.png")
FONT = "Times New Roman"
TITLE_SIZE = Pt(36)
SUBTITLE_SIZE = Pt(24)
BODY_SIZE = Pt(20)


def clone_slide(prs, source_slide):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        new_slide.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)

    return new_slide


def move_slide(prs, old_index, new_index):
    sld_id_lst = prs.slides._sldIdLst
    sld_id = sld_id_lst[old_index]
    sld_id_lst.remove(sld_id)
    sld_id_lst.insert(new_index, sld_id)


def set_text(shape, text, size, *, bold=False, align=PP_ALIGN.LEFT):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = 1.05
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(10)
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = size
        run.font.bold = bold


def remove_pictures(slide):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.element.getparent().remove(shape.element)


def replace_slide_content(slide):
    # The cloned source is the current "Bat dau tran" slide: title, subtitle, body, picture.
    set_text(slide.shapes[0], "GIAO DI\u1ec6N V\u00c0\nTR\u1ea2I NGHI\u1ec6M", TITLE_SIZE, bold=True, align=PP_ALIGN.CENTER)
    set_text(slide.shapes[1], "Nh\u1eadp t\u00ean ng\u01b0\u1eddi ch\u01a1i", SUBTITLE_SIZE, bold=True)
    set_text(
        slide.shapes[2],
        "- Ng\u01b0\u1eddi ch\u01a1i c\u1ea7n nh\u1eadp t\u00ean tr\u01b0\u1edbc khi b\u1eaft \u0111\u1ea7u tr\u1eadn.\n"
        "- T\u00ean \u0111\u01b0\u1ee3c d\u00f9ng \u0111\u1ec3 hi\u1ec3n th\u1ecb k\u1ebft qu\u1ea3 tr\u00ean b\u1ea3ng x\u1ebfp h\u1ea1ng.\n"
        "- N\u1ebfu ch\u01b0a c\u00f3 t\u00ean, h\u1ec7 th\u1ed1ng y\u00eau c\u1ea7u nh\u1eadp v\u00e0 x\u00e1c nh\u1eadn tr\u01b0\u1edbc khi ch\u01a1i.",
        BODY_SIZE,
    )

    old_picture = next((shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE), None)
    left, top, width, height = old_picture.left, old_picture.top, old_picture.width, old_picture.height
    remove_pictures(slide)
    slide.shapes.add_picture(str(IMAGE), left, top, width=width, height=height)


def ensure_fade_transitions(pptx_path):
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    tmp_dir = Path(tempfile.mkdtemp(prefix="pptx_fade_"))
    tmp_pptx = tmp_dir / pptx_path.name

    with ZipFile(pptx_path, "r") as zin, ZipFile(tmp_pptx, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                root = etree.fromstring(data)
                trans = root.find(f"{{{ns}}}transition")
                if trans is None:
                    trans = etree.Element(f"{{{ns}}}transition")
                    clr_map = root.find(f"{{{ns}}}clrMapOvr")
                    if clr_map is not None:
                        root.insert(root.index(clr_map) + 1, trans)
                    else:
                        root.insert(0, trans)
                trans.attrib.pop("advTm", None)
                trans.attrib.pop("advClick", None)
                for child in list(trans):
                    trans.remove(child)
                trans.append(etree.Element(f"{{{ns}}}fade"))
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(item, data)

    shutil.move(str(tmp_pptx), str(pptx_path))
    shutil.rmtree(tmp_dir, ignore_errors=True)


def main():
    prs = Presentation(str(PPTX))

    # Insert after current slide 13, before the existing "Bat dau tran" slide.
    source_slide = prs.slides[14]
    new_slide = clone_slide(prs, source_slide)
    replace_slide_content(new_slide)
    move_slide(prs, len(prs.slides) - 1, 13)

    prs.save(str(PPTX))
    ensure_fade_transitions(PPTX)
    print(PPTX)


if __name__ == "__main__":
    main()
