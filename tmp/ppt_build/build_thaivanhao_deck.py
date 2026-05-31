from pathlib import Path
import shutil

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\Github\Roguelike_Project")
OUTPUT_PPTX = ROOT / "output" / "ppt" / "ThaiVanHao-2121051075_BaoCaoDoAnGame.pptx"
TEMPLATE_PPTX = Path(r"C:\Users\haov8\Downloads\BaoCaoDoAnGame_tren_form_BaoCaoDoAn_final-1.pptx")
ALT_PPTX = ROOT / "output" / "ppt" / "ThaiVanHao-2121051075_BaoCaoDoAnGame_font36_placeholder.pptx"
ASSET_DIR = ROOT / "output" / "ppt" / "assets_thaivanhao"

FONT_NAME = "Times New Roman"
TITLE_SIZE = Pt(36)
BODY_SIZE = Pt(20)
SMALL_SIZE = Pt(18)

TITLE_COLOR = RGBColor(31, 78, 121)
BODY_COLOR = RGBColor(0, 0, 0)
FRAME_COLOR = RGBColor(31, 78, 121)
FRAME_FILL = RGBColor(245, 247, 250)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

TITLE_LEFT = Inches(2.35)
TITLE_TOP = Inches(1.25)
TITLE_WIDTH = Inches(6.80)
TITLE_HEIGHT = Inches(0.65)
LONG_TITLE_TOP = Inches(0.92)
LONG_TITLE_HEIGHT = Inches(1.12)

CONTENT_LEFT = Inches(0.85)
CONTENT_RIGHT = Inches(0.80)
CONTENT_WIDTH = SLIDE_WIDTH - CONTENT_LEFT - CONTENT_RIGHT
SECTION_TOP = Inches(2.22)
TEXT_TOP = Inches(2.72)
LONG_SECTION_TOP = Inches(2.44)
LONG_TEXT_TOP = Inches(2.94)
TEXT_LINE_HEIGHT = Inches(0.78)
TEXT_GAP = Inches(0.06)

DISPLAY_TITLE_OVERRIDES = {
    "LÝ THUYẾT VÀ CÔNG NGHỆ": "LÝ THUYẾT VÀ\nCÔNG NGHỆ",
    "GIAO DIỆN VÀ TRẢI NGHIỆM": "GIAO DIỆN VÀ\nTRẢI NGHIỆM",
}

SLIDES = [
    {
        "kind": "cover",
        "title": "ĐỒ ÁN TỐT NGHIỆP",
        "subtitle": "NGHIÊN CỨU VÀ PHÁT TRIỂN TRÒ CHƠI SINH TỒN 3D",
        "lines": [
            "Sinh viên thực hiện: Thái Văn Hào",
            "Mã sinh viên: 2121051075",
            "Giáo viên hướng dẫn: Phạm Quang Hiển",
        ],
    },
    {
        "kind": "agenda",
        "title": "NỘI DUNG TRÌNH BÀY",
        "lines": [
            "1. GIỚI THIỆU ĐỀ TÀI",
            "2. LÝ THUYẾT VÀ CÔNG NGHỆ",
            "3. PHÂN TÍCH HỆ THỐNG",
            "4. QUY TRÌNH THỰC HIỆN",
            "5. GIAO DIỆN VÀ TRẢI NGHIỆM",
            "6. KẾT LUẬN",
        ],
    },
    {
        "kind": "text",
        "title": "GIỚI THIỆU ĐỀ TÀI",
        "section": "Lý do chọn đề tài:",
        "lines": [
            "Thể loại roguelike sinh tồn đánh quái có tính hấp dẫn, nhịp độ nhanh và giá trị chơi lại cao.",
            "Đề tài phù hợp để áp dụng tổng hợp kiến thức về Unity, C#, gameplay, UI và backend.",
            "Roguelike sinh tồn có vòng chơi ngắn, dễ demo và dễ đánh giá.",
            "Đề tài có tiềm năng mở rộng thành một game hoàn thiện hơn trong tương lai.",
        ],
    },
    {
        "kind": "two_col",
        "title": "GIỚI THIỆU ĐỀ TÀI",
        "left_title": "Mục tiêu",
        "left_lines": [
            "Xây dựng game sinh tồn đánh quái 3D trên Unity.",
            "Hoàn thiện gameplay từ di chuyển, chiến đấu, nhận kinh nghiệm đến chọn buff theo wave.",
            "Xây dựng giao diện và bảng xếp hạng để sản phẩm có thể chạy, build và chơi mượt mà trên các thiết bị khác nhau",
        ],
        "right_title": "Phạm vi",
        "right_lines": [
            "Tập trung vào gameplay chính, enemy, wave, buff, HUD, menu tạm dừng, nhập tên và bảng xếp hạng",
            "Chưa đi sâu vào multiplayer, cốt truyện",
            "Sản phẩm là một bản game cơ bản có thể phát triển thêm nhiều tính năng sau này",
        ],
    },
    {
        "kind": "text_image",
        "title": "LÝ THUYẾT VÀ CÔNG NGHỆ",
        "section": "Tổng quan gameplay",
        "lines": [
            "Người chơi di chuyển, né tránh enemy và tấn công mục tiêu.",
            "Khi tiêu diệt enemy, người chơi nhận kinh nghiệm, lên cấp và chọn buff để tăng sức mạnh.",
            "Độ khó tăng dần theo từng wave và có boss ở các mốc quan trọng.",
            "Trò chơi có HUD và bảng xếp hạng để hỗ trợ trải nghiệm và theo dõi kết quả.",
        ],
        "image_path": ASSET_DIR / "doc_image_34.jpeg",
        "placeholder": "Chèn hình gameplay tổng quan",
    },
    {
        "kind": "tech_grid",
        "title": "LÝ THUYẾT VÀ CÔNG NGHỆ",
        "items": [
            {
                "title": "Unity Engine",
                "placeholder": "Chèn logo / ảnh Unity",
                "role": "Dựng scene, gameplay runtime và UI.",
            },
            {
                "title": "C#",
                "placeholder": "Chèn logo / ảnh C#",
                "role": "Ngôn ngữ lập trình chính",
            },
            {
                "title": "PlayFab",
                "placeholder": "Chèn logo / ảnh PlayFab",
                "role": "Lưu leaderboard và đồng bộ tên người chơi.",
            },
            {
                "title": "Visual Studio Code",
                "placeholder": "Chèn logo / ảnh Visual Studio Code",
                "role": "Môi trường lập trình, xử lý logic",
            },
        ],
    },
    {
        "kind": "two_col",
        "title": "PHÂN TÍCH HỆ THỐNG",
        "section": "Nhóm chức năng chính",
        "left_title": "Gameplay chính",
        "left_lines": [
            "Định danh người chơi, bắt đầu trận đấu",
            "Điều khiển nhân vật, chiến đấu với enemy và vượt qua các wave.",
            "Nhận kinh nghiệm, lên cấp, chọn buff và đổi theme map theo tiến trình.",
        ],
        "right_title": "Giao diện và hệ thống hỗ trợ",
        "right_lines": [
            "Hiển thị HUD, menu tạm dừng, cài đặt và loading hỗ trợ trong trận.",
            "Nhập tên người chơi và hiển thị bảng xếp hạng.",
            "Kết nối các chức năng thành một vòng chơi hoàn chỉnh.",
        ],
    },
    {
        "kind": "two_col",
        "title": "PHÂN TÍCH HỆ THỐNG",
        "section": "Tổ chức project và module code",
        "left_title": "Nhóm gameplay",
        "left_lines": [
            "Xử lý điều khiển nhân vật và chiến đấu.",
            "Quản lý enemy, đạn bắn và wave",
            "Tổ chức dữ liệu để dễ chỉnh sửa và phát triển thêm.",
        ],
        "right_title": "Nhóm hỗ trợ",
        "right_lines": [
            "Quản lý kinh nghiệm, lên cấp và buff.",
            "Hiển thị giao diện và các chức năng hỗ trợ trong game.",
            "Kết nối bảng xếp hạng và lưu điểm người chơi.",
        ],
    },
    {
        "kind": "text_image",
        "title": "QUY TRÌNH THỰC HIỆN",
        "section": "Điều khiển và chiến đấu",
        "lines": [
            "Người chơi di chuyển nhân vật bằng bàn phím và có thể sử dụng dash để né enemy.",
            "Hệ thống tự động tìm mục tiêu gần nhất để thực hiện tấn công.",
            "Cơ chế bắn tự động giúp người chơi tập trung nhiều hơn vào di chuyển và chiến thuật giữ vị trí.",
            "Cách chơi này giúp trận đấu diễn ra nhanh và tạo cảm giác liên tục.",
        ],
        "image_path": ASSET_DIR / "doc_image_35.png",
        "placeholder": "Chèn hình combat",
    },
    {
        "kind": "text_image_grid",
        "title": "QUY TRÌNH THỰC HIỆN",
        "section": "Enemy, boss và wave",
        "lines": [
            "Enemy được chia theo vai trò áp sát, bay hoặc gây áp lực từ xa.",
            "WaveSpawner sinh quái theo cấu hình từng wave.",
            "Boss wave xuất hiện ở các mốc quan trọng để tăng cao trào.",
            "Độ khó tăng dần theo số lượng quái, nhịp spawn và chỉ số.",
        ],
        "images": [
            {"path": ASSET_DIR / "doc_image_37.png", "caption": "Enemy bay"},
            {"path": ASSET_DIR / "doc_image_38.png", "caption": "Enemy đánh xa"},
            {"path": ASSET_DIR / "doc_image_39.png", "caption": "Enemy cận chiến"},
            {"path": ASSET_DIR / "doc_image_40.png", "caption": "Boss / LawaChurl"},
        ],
    },
    {
        "kind": "text_image",
        "title": "QUY TRÌNH THỰC HIỆN",
        "section": "Kinh nghiệm, lên cấp và buff",
        "lines": [
            "Người chơi nhận kinh nghiệm khi tiêu diệt enemy.",
            "Khi đủ kinh nghiệm, nhân vật sẽ lên cấp.",
            "Mỗi lần lên cấp, người chơi được chọn buff để tăng sức mạnh.",
            "Các buff giúp thay đổi cách chơi và hỗ trợ vượt qua các wave khó hơn.",
        ],
        "image_path": ASSET_DIR / "doc_image_01.png",
        "placeholder": "Chèn hình level-up / buff",
    },
    {
        "kind": "text_image",
        "title": "GIAO DIỆN VÀ TRẢI NGHIỆM",
        "section": "HUD, pause, settings, theme map",
        "lines": [
            "HUD hiển thị các thông tin chính như máu, kinh nghiệm, cấp độ và wave hiện tại.",
            "Menu pause và settings hỗ trợ tạm dừng trò chơi và điều chỉnh âm thanh.",
            "Theme map được thay đổi theo từng giai đoạn để tạo cảm giác mới trong quá trình chơi.",
            "Các giao diện được hiển thị đúng theo trạng thái của trận đấu.",
        ],
        "image_path": ASSET_DIR / "doc_image_04.png",
        "placeholder": "Chèn hình pause / settings",
    },
    {
        "kind": "text_image",
        "title": "GIAO DIỆN VÀ TRẢI NGHIỆM",
        "section": "Challenge và bắt đầu trận",
        "lines": [
            "Người chơi tương tác với khu vực challenge để mở panel bắt đầu trận.",
            "Panel hiển thị hướng dẫn điều khiển cơ bản trước khi vào trận.",
            "Sau khi xác nhận bắt đầu, hệ thống khởi tạo wave đầu tiên và chuyển sang trạng thái chiến đấu.",
        ],
        "image_path": ASSET_DIR / "doc_image_47.png",
        "placeholder": "Chèn hình challenge bắt đầu",
    },
    {
        "kind": "text_image",
        "title": "GIAO DIỆN VÀ TRẢI NGHIỆM",
        "section": "Leaderboard",
        "lines": [
            "Điểm số của người chơi được gửi lên hệ thống PlayFab sau khi kết thúc trận đấu.",
            "Bảng xếp hạng hiển thị danh sách người chơi theo điểm số từ cao xuống thấp.",
            "Người chơi có thể theo dõi thứ hạng của mình trực tiếp trên giao diện.",
        ],
        "image_path": ASSET_DIR / "doc_image_02.png",
        "placeholder": "Chèn hình leaderboard",
    },
    {
        "kind": "text",
        "title": "KẾT LUẬN",
        "section": "Kiểm thử và kết quả đạt được",
        "lines": [
            "Đề tài đã kiểm thử các chức năng chính như vào trận, chiến đấu, lên cấp, tạm dừng và bảng xếp hạng.",
            "Kết quả đạt được là xây dựng được một game sinh tồn 3D có thể chạy và trình bày trên thực tế.",
            "Sản phẩm đã thể hiện rõ phần gameplay, giao diện và kết nối bảng xếp hạng trực tuyến.",
        ],
    },
    {
        "kind": "text",
        "title": "KẾT LUẬN",
        "section": "Hạn chế và hướng phát triển",
        "lines": [
            "Số lượng enemy, buff và nội dung trong game hiện vẫn còn chưa nhiều.",
            "Một số phần giao diện và hiệu năng vẫn cần tiếp tục hoàn thiện.",
            "Trong thời gian tới, đề tài có thể mở rộng thêm map, boss, buff và các chỉ số thống kê.",
        ],
    },
    {
        "kind": "thanks",
        "title": "EM XIN CẢM ƠN THẦY CÔ ĐÃ LẮNG NGHE",
    },
]


def clear_slide(slide):
    for shape in list(slide.shapes):
        element = shape.element if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER else shape._element
        element.getparent().remove(element)


def ensure_slide_count(prs, count):
    while len(prs.slides) < count:
        prs.slides.add_slide(prs.slide_layouts[1])
    while len(prs.slides) > count:
        slide_id = prs.slides._sldIdLst[-1]
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def set_shape_name(shape, name):
    for element in shape._element.iter():
        if element.tag.endswith("}cNvPr"):
            element.set("name", name)
            break


def add_run(paragraph, text, size, bold=False, color=BODY_COLOR):
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    size,
    *,
    bold=False,
    align=PP_ALIGN.LEFT,
    color=BODY_COLOR,
    vertical_anchor=MSO_ANCHOR.TOP,
    margin=0.02,
    name=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = vertical_anchor

    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = 1.05
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    add_run(paragraph, text, size, bold=bold, color=color)

    if name:
        set_shape_name(box, name)
    return box


def add_title(slide, text):
    display_text = DISPLAY_TITLE_OVERRIDES.get(text, text)
    title_top = LONG_TITLE_TOP if text in DISPLAY_TITLE_OVERRIDES else TITLE_TOP
    title_height = LONG_TITLE_HEIGHT if text in DISPLAY_TITLE_OVERRIDES else TITLE_HEIGHT
    add_text_box(
        slide,
        TITLE_LEFT,
        title_top,
        TITLE_WIDTH,
        title_height,
        display_text,
        TITLE_SIZE,
        bold=True,
        align=PP_ALIGN.CENTER,
        color=TITLE_COLOR,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def section_top_for_title(title):
    return LONG_SECTION_TOP if title in DISPLAY_TITLE_OVERRIDES else SECTION_TOP


def text_top_for_title(title):
    return LONG_TEXT_TOP if title in DISPLAY_TITLE_OVERRIDES else TEXT_TOP


def add_section_heading(slide, text, *, title):
    add_text_box(
        slide,
        CONTENT_LEFT,
        section_top_for_title(title),
        CONTENT_WIDTH,
        Inches(0.35),
        text,
        BODY_SIZE,
        bold=True,
        color=BODY_COLOR,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_bullet_lines(slide, left, top, width, lines, *, start_order=10):
    y = top
    order = start_order
    for index, line in enumerate(lines, start=1):
        text = f"- {line}"
        add_text_box(
            slide,
            left,
            y,
            width,
            TEXT_LINE_HEIGHT,
            text,
            BODY_SIZE,
            name=f"anim-appear-{order:02d}-line-{index}",
        )
        y += TEXT_LINE_HEIGHT + TEXT_GAP
        order += 10


def add_center_lines(slide, top, lines, *, start_order=10):
    y = top
    order = start_order
    for index, line in enumerate(lines, start=1):
        add_text_box(
            slide,
            CONTENT_LEFT,
            y,
            CONTENT_WIDTH,
            Inches(0.36),
            line,
            BODY_SIZE,
            align=PP_ALIGN.CENTER,
            name=f"anim-appear-{order:02d}-agenda-{index}",
        )
        y += Inches(0.48)
        order += 10


def add_frame(slide, left, top, width, height):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = FRAME_FILL
    rect.line.color.rgb = FRAME_COLOR
    rect.line.width = Pt(1.1)
    return rect


def add_placeholder(slide, left, top, width, height, caption):
    rect = add_frame(slide, left, top, width, height)
    frame = rect.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.line_spacing = 1.0
    add_run(paragraph, caption, BODY_SIZE, color=FRAME_COLOR)
    return rect


def add_picture_in_frame(slide, image_path, left, top, width, height, *, name=None):
    add_frame(slide, left, top, width, height)
    inset = Inches(0.06)
    inner_left = left + inset
    inner_top = top + inset
    inner_width = width - inset * 2
    inner_height = height - inset * 2

    with Image.open(image_path) as image:
        image_width, image_height = image.size

    image_ratio = image_width / image_height
    frame_ratio = inner_width / inner_height

    if image_ratio >= frame_ratio:
        final_width = inner_width
        final_height = round(inner_width / image_ratio)
        final_left = inner_left
        final_top = inner_top + round((inner_height - final_height) / 2)
    else:
        final_height = inner_height
        final_width = round(inner_height * image_ratio)
        final_left = inner_left + round((inner_width - final_width) / 2)
        final_top = inner_top

    picture = slide.shapes.add_picture(
        str(image_path),
        final_left,
        final_top,
        width=final_width,
        height=final_height,
    )
    picture.line.color.rgb = FRAME_COLOR
    picture.line.width = Pt(0.7)
    if name:
        set_shape_name(picture, name)
    return picture


def add_media(slide, left, top, width, height, *, image_path=None, placeholder="", anim_name=None):
    if image_path and Path(image_path).exists():
        return add_picture_in_frame(slide, Path(image_path), left, top, width, height, name=anim_name)
    shape = add_placeholder(slide, left, top, width, height, placeholder)
    if anim_name:
        set_shape_name(shape, anim_name)
    return shape


def build_cover(slide, spec):
    add_text_box(
        slide,
        Inches(2.70),
        Inches(1.75),
        Inches(8.00),
        Inches(0.45),
        spec["title"],
        Pt(24),
        bold=True,
        align=PP_ALIGN.CENTER,
        color=TITLE_COLOR,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide,
        Inches(1.55),
        Inches(2.45),
        Inches(10.20),
        Inches(1.10),
        spec["subtitle"],
        TITLE_SIZE,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    y = Inches(4.25)
    for line in spec["lines"]:
        add_text_box(
            slide,
            Inches(2.10),
            y,
            Inches(9.10),
            Inches(0.35),
            line,
            BODY_SIZE,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        y += Inches(0.42)


def build_agenda(slide, spec):
    add_title(slide, spec["title"])
    add_center_lines(slide, Inches(2.70), spec["lines"])


def build_text(slide, spec):
    add_title(slide, spec["title"])
    add_section_heading(slide, spec["section"], title=spec["title"])
    add_bullet_lines(slide, CONTENT_LEFT, text_top_for_title(spec["title"]), Inches(11.00), spec["lines"])


def build_two_col(slide, spec):
    add_title(slide, spec["title"])
    if spec.get("section"):
        add_section_heading(slide, spec["section"], title=spec["title"])
        title_y = Inches(2.68)
    else:
        title_y = Inches(2.45)

    gap = Inches(0.45)
    col_width = (CONTENT_WIDTH - gap) / 2
    left_x = CONTENT_LEFT
    right_x = CONTENT_LEFT + col_width + gap

    add_text_box(
        slide,
        left_x,
        title_y,
        col_width,
        Inches(0.35),
        spec["left_title"],
        BODY_SIZE,
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide,
        right_x,
        title_y,
        col_width,
        Inches(0.35),
        spec["right_title"],
        BODY_SIZE,
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_bullet_lines(slide, left_x, title_y + Inches(0.45), col_width, spec["left_lines"], start_order=10)
    add_bullet_lines(slide, right_x, title_y + Inches(0.45), col_width, spec["right_lines"], start_order=40)


def build_text_image(slide, spec):
    add_title(slide, spec["title"])
    add_section_heading(slide, spec["section"], title=spec["title"])
    media_top = text_top_for_title(spec["title"])
    add_media(
        slide,
        Inches(7.15),
        media_top,
        Inches(5.10),
        Inches(3.45),
        image_path=spec.get("image_path"),
        placeholder=spec["placeholder"],
        anim_name="anim-fade-10-media",
    )
    add_bullet_lines(slide, CONTENT_LEFT, media_top, Inches(5.70), spec["lines"], start_order=20)


def build_text_image_grid(slide, spec):
    add_title(slide, spec["title"])
    add_section_heading(slide, spec["section"], title=spec["title"])
    media_top = text_top_for_title(spec["title"])
    add_bullet_lines(slide, CONTENT_LEFT, media_top, Inches(5.55), spec["lines"], start_order=50)

    grid_left = Inches(7.10)
    grid_top = media_top
    col_gap = Inches(0.18)
    row_gap = Inches(0.14)
    image_width = Inches(2.42)
    image_height = Inches(1.22)
    caption_height = Inches(0.28)

    positions = [
        (grid_left, grid_top),
        (grid_left + image_width + col_gap, grid_top),
        (grid_left, grid_top + image_height + caption_height + row_gap),
        (grid_left + image_width + col_gap, grid_top + image_height + caption_height + row_gap),
    ]

    for index, (item, (x, y)) in enumerate(zip(spec["images"], positions), start=1):
        add_media(
            slide,
            x,
            y,
            image_width,
            image_height,
            image_path=item["path"],
            placeholder="Chèn hình enemy",
            anim_name=f"anim-fade-{index * 10:02d}-enemy-{index}",
        )
        add_text_box(
            slide,
            x,
            y + image_height + Inches(0.03),
            image_width,
            caption_height,
            item["caption"],
            BODY_SIZE,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
            name=f"anim-appear-{index * 10 + 1:02d}-enemy-cap-{index}",
        )


def build_tech_grid(slide, spec):
    add_title(slide, spec["title"])
    top_y = Inches(2.58) if spec["title"] in DISPLAY_TITLE_OVERRIDES else Inches(2.40)
    row_gap = Inches(0.38)
    col_gap = Inches(0.40)
    cell_width = (CONTENT_WIDTH - col_gap) / 2
    left_x = CONTENT_LEFT
    right_x = CONTENT_LEFT + cell_width + col_gap
    positions = [
        (left_x, top_y),
        (right_x, top_y),
        (left_x, top_y + Inches(1.92) + row_gap),
        (right_x, top_y + Inches(1.92) + row_gap),
    ]

    for item, (cell_x, cell_y) in zip(spec["items"], positions):
        add_text_box(
            slide,
            cell_x,
            cell_y,
            cell_width,
            Inches(0.35),
            item["title"],
            BODY_SIZE,
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        add_placeholder(
            slide,
            cell_x + Inches(0.04),
            cell_y + Inches(0.42),
            cell_width - Inches(0.08),
            Inches(0.92),
            item["placeholder"],
        )
        add_text_box(
            slide,
            cell_x + Inches(0.02),
            cell_y + Inches(1.42),
            cell_width - Inches(0.04),
            Inches(0.40),
            item["role"],
            BODY_SIZE,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )


def build_demo_grid(slide, spec):
    add_title(slide, spec["title"])
    add_section_heading(slide, spec["section"], title=spec["title"])

    image_width = Inches(3.45)
    image_height = Inches(2.35)
    image_top = text_top_for_title(spec["title"])
    caption_top = image_top + image_height + Inches(0.11)
    left_positions = [Inches(0.95), Inches(4.95), Inches(8.95)]

    for index, (item, x) in enumerate(zip(spec["images"], left_positions), start=1):
        add_media(
            slide,
            x,
            image_top,
            image_width,
            image_height,
            image_path=item["path"],
            placeholder="Chèn hình demo",
            anim_name=f"anim-fade-{index * 10:02d}-demo-{index}",
        )
        add_text_box(
            slide,
            x,
            caption_top,
            image_width,
            Inches(0.35),
            item["caption"],
            BODY_SIZE,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    add_text_box(
        slide,
        CONTENT_LEFT,
        caption_top + Inches(0.62),
        CONTENT_WIDTH,
        Inches(0.55),
        spec["note"],
        BODY_SIZE,
        align=PP_ALIGN.CENTER,
        name="anim-appear-40-demo-note",
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def build_closing(slide, spec):
    add_title(slide, spec["title"])
    add_section_heading(slide, spec["section"], title=spec["title"])
    add_bullet_lines(slide, CONTENT_LEFT, text_top_for_title(spec["title"]), Inches(11.00), spec["lines"])
    add_text_box(
        slide,
        Inches(1.40),
        Inches(5.95),
        Inches(10.60),
        Inches(0.42),
        spec["thanks"],
        BODY_SIZE,
        bold=True,
        align=PP_ALIGN.CENTER,
        color=TITLE_COLOR,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def build_thanks(slide, spec):
    add_text_box(
        slide,
        Inches(1.15),
        Inches(3.05),
        Inches(11.10),
        Inches(0.85),
        spec["title"],
        TITLE_SIZE,
        bold=True,
        align=PP_ALIGN.CENTER,
        color=TITLE_COLOR,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


BUILDERS = {
    "cover": build_cover,
    "agenda": build_agenda,
    "text": build_text,
    "two_col": build_two_col,
    "text_image": build_text_image,
    "text_image_grid": build_text_image_grid,
    "tech_grid": build_tech_grid,
    "demo_grid": build_demo_grid,
    "closing": build_closing,
    "thanks": build_thanks,
}


def verify(prs):
    slides = list(prs.slides)
    if len(slides) != 17:
        raise RuntimeError(f"Expected 17 slides, got {len(slides)}")

    issues = []
    for slide_index, slide in enumerate(slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text.strip()
                    if not text:
                        continue
                    if run.font.name != FONT_NAME:
                        issues.append((slide_index, text, run.font.name))
                    size = round(run.font.size.pt, 2) if run.font.size else None
                    if slide_index == 1:
                        if size not in {20.0, 24.0, 36.0}:
                            issues.append((slide_index, text, size))
                    elif slide_index == 17:
                        if size not in {20.0, 36.0}:
                            issues.append((slide_index, text, size))
                    else:
                        if size not in {20.0, 36.0}:
                            issues.append((slide_index, text, size))
                    lowered = text.lower()
                    if any(word in lowered for word in ["bộ môn tin học kinh tế", "diagram", "flowchart", "sơ đồ"]):
                        issues.append((slide_index, text, "unexpected-text"))

    if issues:
        raise RuntimeError(f"Verification issues: {issues[:12]}")


def build():
    if not TEMPLATE_PPTX.exists():
        raise FileNotFoundError(f"Missing source deck: {TEMPLATE_PPTX}")

    prs = Presentation(str(TEMPLATE_PPTX))
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    ensure_slide_count(prs, len(SLIDES))

    slides = list(prs.slides)
    for slide, spec in zip(slides, SLIDES):
        clear_slide(slide)
        BUILDERS[spec["kind"]](slide, spec)

    verify(prs)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    temp_path = OUTPUT_PPTX.with_suffix(".tmp.pptx")
    prs.save(str(temp_path))

    try:
        shutil.copyfile(temp_path, OUTPUT_PPTX)
        temp_path.unlink(missing_ok=True)
        return OUTPUT_PPTX
    except PermissionError:
        shutil.copyfile(temp_path, ALT_PPTX)
        temp_path.unlink(missing_ok=True)
        return ALT_PPTX


if __name__ == "__main__":
    print(build())
