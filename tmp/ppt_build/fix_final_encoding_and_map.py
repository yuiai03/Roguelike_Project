from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


PPTX = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")
FONT = "Times New Roman"
BODY_SIZE = Pt(20)


def set_paragraph(paragraph, text):
    paragraph.clear()
    paragraph.line_spacing = 1.05
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(10)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = BODY_SIZE


def main():
    prs = Presentation(str(PPTX))
    set_paragraph(
        prs.slides[8].shapes[3].text_frame.paragraphs[1],
        "- Qu\u1ea3n l\u00fd kinh nghi\u1ec7m, l\u00ean c\u1ea5p, n\u00e2ng c\u1ea5p v\u00e0 b\u1ea3n \u0111\u1ed3.",
    )
    set_paragraph(
        prs.slides[18].shapes[1].text_frame.paragraphs[3],
        "- Trong th\u1eddi gian t\u1edbi, \u0111\u1ec1 t\u00e0i c\u00f3 th\u1ec3 m\u1edf r\u1ed9ng th\u00eam lo\u1ea1i b\u1ea3n \u0111\u1ed3, lo\u1ea1i qu\u00e1i, nhi\u1ec1u m\u00f4 h\u00ecnh nh\u00e2n v\u1eadt, n\u00e2ng c\u1ea5p, giao di\u1ec7n v\u00e0 c\u00e1c t\u00ednh n\u0103ng kh\u00e1c th\u00fa v\u1ecb h\u01a1n.",
    )
    prs.save(str(PPTX))
    print(PPTX)


if __name__ == "__main__":
    main()
