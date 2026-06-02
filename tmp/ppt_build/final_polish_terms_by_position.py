from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


PPTX = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")
FONT = "Times New Roman"
BODY_SIZE = Pt(20)
SUBTITLE_SIZE = Pt(24)


def set_paragraph(paragraph, text, size=BODY_SIZE, bold=False):
    paragraph.clear()
    paragraph.line_spacing = 1.05
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(10)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold


def main():
    prs = Presentation(str(PPTX))

    # Slide 3
    set_paragraph(
        prs.slides[2].shapes[2].text_frame.paragraphs[1],
        "- \u0110\u1ec1 t\u00e0i ph\u00f9 h\u1ee3p \u0111\u1ec3 \u00e1p d\u1ee5ng nh\u1eefng ki\u1ebfn th\u1ee9c v\u1ec1 Unity, C#, l\u1ed1i ch\u01a1i, giao di\u1ec7n v\u00e0 k\u1ebft n\u1ed1i d\u1eef li\u1ec7u.",
    )

    # Slide 4
    set_paragraph(
        prs.slides[3].shapes[4].text_frame.paragraphs[1],
        "- Ch\u01b0a \u0111i s\u00e2u v\u00e0o ch\u01a1i nhi\u1ec1u ng\u01b0\u1eddi v\u00e0 c\u1ed1t truy\u1ec7n",
    )

    # Slide 9
    set_paragraph(
        prs.slides[8].shapes[1].text_frame.paragraphs[0],
        "T\u1ed5 ch\u1ee9c m\u00e3 ngu\u1ed3n v\u00e0 nh\u00f3m ch\u1ee9c n\u0103ng",
        size=SUBTITLE_SIZE,
        bold=True,
    )
    set_paragraph(
        prs.slides[8].shapes[3].text_frame.paragraphs[0],
        "Nh\u00f3m h\u1ed7 tr\u1ee3",
        bold=True,
    )
    set_paragraph(
        prs.slides[8].shapes[3].text_frame.paragraphs[1],
        "- Qu\u1ea3n l\u00fd kinh nghi\u1ec7m, l\u00ean c\u1ea5p, n\u00e2ng c\u1ea5p v\u00e0 b\u1ea3n \u0111\u1ed3.",
    )

    # Slide 11
    set_paragraph(
        prs.slides[10].shapes[1].text_frame.paragraphs[0],
        "Qu\u00e1i, boss v\u00e0 c\u00e1c \u0111\u1ee3t t\u1ea5n c\u00f4ng",
        size=SUBTITLE_SIZE,
        bold=True,
    )
    set_paragraph(
        prs.slides[10].shapes[2].text_frame.paragraphs[1],
        "- C\u00e1c lo\u1ea1i qu\u00e1i xu\u1ea5t hi\u1ec7n theo t\u1eebng \u0111\u1ee3t \u0111\u1ec3 t\u1ea1o \u00e1p l\u1ef1c cho ng\u01b0\u1eddi ch\u01a1i.",
    )
    set_paragraph(
        prs.slides[10].shapes[2].text_frame.paragraphs[4],
        "- Ho\u00e0n th\u00e0nh \u0111\u1ee3t c\u00f3 boss s\u1ebd chuy\u1ec3n sang b\u1ea3n \u0111\u1ed3 m\u1edbi.",
    )

    # Slide 13
    set_paragraph(
        prs.slides[12].shapes[2].text_frame.paragraphs[0],
        "Giao di\u1ec7n trong tr\u1eadn, t\u1ea1m d\u1eebng v\u00e0 b\u1ea3n \u0111\u1ed3",
        size=SUBTITLE_SIZE,
        bold=True,
    )
    set_paragraph(
        prs.slides[12].shapes[2].text_frame.paragraphs[1],
        "- Giao di\u1ec7n trong tr\u1eadn hi\u1ec3n th\u1ecb m\u00e1u, kinh nghi\u1ec7m, c\u1ea5p \u0111\u1ed9 v\u00e0 \u0111\u1ee3t qu\u00e1i hi\u1ec7n t\u1ea1i.",
    )
    set_paragraph(
        prs.slides[12].shapes[2].text_frame.paragraphs[3],
        "- B\u1ea3n \u0111\u1ed3 thay \u0111\u1ed5i theo t\u1eebng giai \u0111o\u1ea1n \u0111\u1ec3 t\u1ea1o c\u1ea3m gi\u00e1c m\u1edbi trong qu\u00e1 tr\u00ecnh ch\u01a1i.",
    )

    prs.save(str(PPTX))
    print(PPTX)


if __name__ == "__main__":
    main()
