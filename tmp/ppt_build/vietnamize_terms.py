from pathlib import Path

from pptx import Presentation


PPTX = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")

REPLACEMENTS = [
    ("gameplay", "l\u1ed1i ch\u01a1i"),
    ("Gameplay", "L\u1ed1i ch\u01a1i"),
    ("Player", "Nh\u00e2n v\u1eadt"),
    ("player", "nh\u00e2n v\u1eadt"),
    ("enemy", "qu\u00e1i"),
    ("Enemy", "Qu\u00e1i"),
    ("wave", "\u0111\u1ee3t qu\u00e1i"),
    ("Wave", "\u0110\u1ee3t qu\u00e1i"),
    ("buff", "n\u00e2ng c\u1ea5p"),
    ("Buff", "N\u00e2ng c\u1ea5p"),
    ("HUD", "giao di\u1ec7n th\u00f4ng tin trong tr\u1eadn"),
    ("UI", "giao di\u1ec7n"),
    ("Leaderboard", "B\u1ea3ng x\u1ebfp h\u1ea1ng"),
    ("leaderboard", "b\u1ea3ng x\u1ebfp h\u1ea1ng"),
    ("challenge", "khu v\u1ef1c b\u1eaft \u0111\u1ea7u tr\u1eadn"),
    ("Challenge", "Khu v\u1ef1c b\u1eaft \u0111\u1ea7u tr\u1eadn"),
    ("panel", "b\u1ea3ng h\u01b0\u1edbng d\u1eabn"),
    ("Panel", "B\u1ea3ng h\u01b0\u1edbng d\u1eabn"),
    ("pause", "t\u1ea1m d\u1eebng"),
    ("Pause", "T\u1ea1m d\u1eebng"),
    ("settings", "c\u00e0i \u0111\u1eb7t"),
    ("Settings", "C\u00e0i \u0111\u1eb7t"),
    ("loading", "m\u00e0n h\u00ecnh ch\u1edd"),
    ("Loading", "M\u00e0n h\u00ecnh ch\u1edd"),
    ("Theme map", "B\u1ea3n \u0111\u1ed3 theo giai \u0111o\u1ea1n"),
    ("theme map", "b\u1ea3n \u0111\u1ed3 theo giai \u0111o\u1ea1n"),
    ("map", "b\u1ea3n \u0111\u1ed3"),
    ("Map", "B\u1ea3n \u0111\u1ed3"),
    ("project", "d\u1ef1 \u00e1n"),
    ("Project", "D\u1ef1 \u00e1n"),
    ("module code", "nh\u00f3m ch\u1ee9c n\u0103ng"),
    ("Module code", "Nh\u00f3m ch\u1ee9c n\u0103ng"),
    ("module", "nh\u00f3m ch\u1ee9c n\u0103ng"),
    ("Module", "Nh\u00f3m ch\u1ee9c n\u0103ng"),
    ("window", "Windows"),
    ("model", "m\u00f4 h\u00ecnh"),
    ("Model", "M\u00f4 h\u00ecnh"),
]

EXACT_TEXT = {
    "T\u1ed5 ch\u1ee9c d\u1ef1 \u00e1n v\u00e0 nh\u00f3m ch\u1ee9c n\u0103ng": "T\u1ed5 ch\u1ee9c m\u00e3 ngu\u1ed3n v\u00e0 nh\u00f3m ch\u1ee9c n\u0103ng",
    "Enemy, boss v\u00e0 c\u00e1c \u0111\u1ee3t qu\u00e1i": "Qu\u00e1i, boss v\u00e0 c\u00e1c \u0111\u1ee3t t\u1ea5n c\u00f4ng",
    "Kinh nghi\u1ec7m, l\u00ean c\u1ea5p v\u00e0 n\u00e2ng c\u1ea5p": "Kinh nghi\u1ec7m, l\u00ean c\u1ea5p v\u00e0 n\u00e2ng c\u1ea5p",
}

PHRASE_FIXES = [
    ("n\u00e2ng c\u1ea5p theo \u0111\u1ee3t qu\u00e1i", "n\u00e2ng c\u1ea5p theo t\u1eebng \u0111\u1ee3t qu\u00e1i"),
    ("ng\u01b0\u1eddi ch\u01a1i, qu\u00e1i, \u0111\u1ee3t qu\u00e1i, n\u00e2ng c\u1ea5p, giao di\u1ec7n", "ng\u01b0\u1eddi ch\u01a1i, qu\u00e1i, \u0111\u1ee3t qu\u00e1i, n\u00e2ng c\u1ea5p v\u00e0 giao di\u1ec7n"),
    ("Qu\u00e1i bay", "Qu\u00e1i bay"),
    ("Qu\u00e1i \u0111\u00e1nh xa", "Qu\u00e1i \u0111\u00e1nh xa"),
    ("Qu\u00e1i c\u1eadn chi\u1ebfn", "Qu\u00e1i c\u1eadn chi\u1ebfn"),
    ("\u0111\u1ee3t qu\u00e1i boss", "\u0111\u1ee3t c\u00f3 boss"),
    ("Windows.", "Windows."),
    ("tr\u00ean Windows", "tr\u00ean Windows"),
]


def replace_text(text):
    result = text
    for old, new in REPLACEMENTS:
        result = result.replace(old, new)
    for old, new in PHRASE_FIXES:
        result = result.replace(old, new)
    return EXACT_TEXT.get(result, result)


def main():
    prs = Presentation(str(PPTX))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        run.text = replace_text(run.text)

    prs.save(str(PPTX))
    print(PPTX)


if __name__ == "__main__":
    main()
