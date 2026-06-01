from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor


SOURCE = Path(r"C:\Users\haov8\Downloads\ThaiVanHao-2121051075-DATHKT.pptx")
FALLBACK = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")
FONT = "Times New Roman"
TITLE_SIZE = Pt(36)
TOOL_SIZE = Pt(24)
DESC_SIZE = Pt(20)
BLUE = RGBColor(31, 78, 121)


def set_text(shape, text, size, *, bold=False, align=PP_ALIGN.LEFT, color=None):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = 1.05
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def main():
    prs = Presentation(str(SOURCE))
    slide = prs.slides[6]

    set_text(
        slide.shapes[0],
        "L\u00dd THUY\u1ebeT V\u00c0\nC\u00d4NG NGH\u1ec6",
        TITLE_SIZE,
        bold=True,
        align=PP_ALIGN.CENTER,
        color=BLUE,
    )
    set_text(slide.shapes[1], "Git/GitHub", TOOL_SIZE, bold=True)
    set_text(
        slide.shapes[2],
        "Qu\u1ea3n l\u00fd phi\u00ean b\u1ea3n, l\u01b0u tr\u1eef project v\u00e0 theo d\u00f5i thay \u0111\u1ed5i m\u00e3 ngu\u1ed3n.",
        DESC_SIZE,
    )
    set_text(slide.shapes[3], "Sketchfab", TOOL_SIZE, bold=True)
    set_text(
        slide.shapes[4],
        "Tham kh\u1ea3o v\u00e0 l\u1ea5y model 3D ph\u00f9 h\u1ee3p \u0111\u1ec3 s\u1eed d\u1ee5ng trong game.",
        DESC_SIZE,
    )
    set_text(slide.shapes[5], "Blender", TOOL_SIZE, bold=True)
    set_text(
        slide.shapes[6],
        "Ch\u1ec9nh s\u1eeda model, t\u1ed1i \u01b0u nh\u00e2n v\u1eadt v\u00e0 h\u1ed7 tr\u1ee3 l\u00e0m ho\u1ea1t \u1ea3nh.",
        DESC_SIZE,
    )
    set_text(slide.shapes[11], "Photoshop", TOOL_SIZE, bold=True)
    set_text(
        slide.shapes[12],
        "Ch\u1ec9nh s\u1eeda h\u00ecnh \u1ea3nh, texture v\u00e0 m\u1ed9t s\u1ed1 t\u00e0i nguy\u00ean giao di\u1ec7n.",
        DESC_SIZE,
    )

    subtitle = "C\u00f4ng c\u1ee5 h\u1ed7 tr\u1ee3 ph\u00e1t tri\u1ec3n"
    subtitle_shape = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if text == subtitle or "C?ng c?" in text:
            subtitle_shape = shape
            break
    if subtitle_shape is not None:
        set_text(subtitle_shape, subtitle, TOOL_SIZE, bold=True)

    try:
        prs.save(str(SOURCE))
        print(SOURCE)
    except PermissionError:
        FALLBACK.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(FALLBACK))
        print(FALLBACK)


if __name__ == "__main__":
    main()
