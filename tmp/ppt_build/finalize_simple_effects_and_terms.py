from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
import shutil

from lxml import etree
from pptx import Presentation
from pptx.util import Pt


PPTX = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")
TMP_DIR = Path(r"E:\Github\Roguelike_Project\tmp\ppt_build\simple_effects_tmp")
FONT = "Times New Roman"
BODY_SIZE = Pt(20)
SUBTITLE_SIZE = Pt(24)


def set_paragraph(paragraph, text, size=BODY_SIZE, bold=False):
    paragraph.clear()
    paragraph.line_spacing = 1.05
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(10)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold


def polish_terms():
    prs = Presentation(str(PPTX))

    set_paragraph(
        prs.slides[2].shapes[2].text_frame.paragraphs[1],
        "- \u0110\u1ec1 t\u00e0i ph\u00f9 h\u1ee3p \u0111\u1ec3 \u00e1p d\u1ee5ng nh\u1eefng ki\u1ebfn th\u1ee9c v\u1ec1 Unity, C#, l\u1ed1i ch\u01a1i, giao di\u1ec7n v\u00e0 k\u1ebft n\u1ed1i d\u1eef li\u1ec7u.",
    )
    set_paragraph(
        prs.slides[3].shapes[3].text_frame.paragraphs[1],
        "- Ho\u00e0n thi\u1ec7n l\u1ed1i ch\u01a1i t\u1eeb di chuy\u1ec3n, chi\u1ebfn \u0111\u1ea5u, nh\u1eadn kinh nghi\u1ec7m \u0111\u1ebfn ch\u1ecdn n\u00e2ng c\u1ea5p theo \u0111\u1ee3t qu\u00e1i.",
    )
    set_paragraph(prs.slides[4].shapes[1].text_frame.paragraphs[0], "T\u1ed5ng quan l\u1ed1i ch\u01a1i", size=SUBTITLE_SIZE, bold=True)
    set_paragraph(prs.slides[5].shapes[2].text_frame.paragraphs[0], "D\u1ef1ng c\u1ea3nh, l\u1ed1i ch\u01a1i v\u00e0 giao di\u1ec7n.")
    set_paragraph(prs.slides[7].shapes[2].text_frame.paragraphs[0], "L\u1ed1i ch\u01a1i", bold=True)
    set_paragraph(
        prs.slides[7].shapes[2].text_frame.paragraphs[3],
        "- Nh\u1eadn kinh nghi\u1ec7m, l\u00ean c\u1ea5p, ch\u1ecdn n\u00e2ng c\u1ea5p v\u00e0 \u0111\u1ed5i b\u1ea3n \u0111\u1ed3 theo ti\u1ebfn tr\u00ecnh.",
    )
    set_paragraph(
        prs.slides[7].shapes[3].text_frame.paragraphs[1],
        "- Hi\u1ec3n th\u1ecb giao di\u1ec7n th\u00f4ng tin trong tr\u1eadn, menu t\u1ea1m d\u1eebng, c\u00e0i \u0111\u1eb7t v\u00e0 m\u00e0n h\u00ecnh ch\u1edd h\u1ed7 tr\u1ee3 trong tr\u1eadn.",
    )
    set_paragraph(prs.slides[8].shapes[2].text_frame.paragraphs[0], "Nh\u00f3m l\u1ed1i ch\u01a1i", bold=True)
    set_paragraph(
        prs.slides[17].shapes[1].text_frame.paragraphs[3],
        "- S\u1ea3n ph\u1ea9m \u0111\u00e3 th\u1ec3 hi\u1ec7n r\u00f5 ph\u1ea7n l\u1ed1i ch\u01a1i, giao di\u1ec7n v\u00e0 k\u1ebft n\u1ed1i b\u1ea3ng x\u1ebfp h\u1ea1ng tr\u1ef1c tuy\u1ebfn.",
    )

    prs.save(str(PPTX))


def minimal_timing(ns):
    timing = etree.Element(f"{{{ns}}}timing")
    tn_lst = etree.SubElement(timing, f"{{{ns}}}tnLst")
    par = etree.SubElement(tn_lst, f"{{{ns}}}par")
    c_tn = etree.SubElement(par, f"{{{ns}}}cTn")
    c_tn.set("id", "1")
    c_tn.set("dur", "indef")
    c_tn.set("restart", "never")
    c_tn.set("nodeType", "tmRoot")
    return timing


def normalize_effects():
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    if TMP_DIR.exists():
        shutil.rmtree(TMP_DIR)
    TMP_DIR.mkdir(parents=True, exist_ok=True)
    tmp_pptx = TMP_DIR / PPTX.name

    with ZipFile(PPTX, "r") as zin, ZipFile(tmp_pptx, "w", ZIP_DEFLATED) as zout:
        slide_count = len([name for name in zin.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                slide_num = int(item.filename.rsplit("slide", 1)[1].split(".xml")[0])
                root = etree.fromstring(data)

                trans = root.find(f"{{{ns}}}transition")
                if trans is None:
                    trans = etree.Element(f"{{{ns}}}transition")
                    clr_map = root.find(f"{{{ns}}}clrMapOvr")
                    if clr_map is not None:
                        root.insert(root.index(clr_map) + 1, trans)
                    else:
                        root.insert(0, trans)
                trans.attrib.pop("advTm", None)
                trans.attrib.pop("advClick", None)
                for child in list(trans):
                    trans.remove(child)
                trans.append(etree.Element(f"{{{ns}}}fade"))

                timing = root.find(f"{{{ns}}}timing")
                if slide_num in (1, slide_count):
                    if timing is not None:
                        root.remove(timing)
                elif timing is None:
                    root.append(minimal_timing(ns))

                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(item, data)

    shutil.move(str(tmp_pptx), str(PPTX))
    shutil.rmtree(TMP_DIR, ignore_errors=True)


def main():
    polish_terms()
    normalize_effects()
    print(PPTX)


if __name__ == "__main__":
    main()
