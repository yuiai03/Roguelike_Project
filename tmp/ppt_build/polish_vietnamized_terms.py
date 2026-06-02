from pathlib import Path

from pptx import Presentation


PPTX = Path(r"E:\Github\Roguelike_Project\output\ppt\ThaiVanHao-2121051075-DATHKT-fixed.pptx")

REPLACEMENTS = {
    "l\u1ed1i ch\u01a1i, giao di\u1ec7n v\u00e0 backend": "l\u1ed1i ch\u01a1i, giao di\u1ec7n v\u00e0 k\u1ebft n\u1ed1i d\u1eef li\u1ec7u",
    "- Ch\u01b0a \u0111i s\u00e2u v\u00e0o multinh\u00e2n v\u1eadt, c\u1ed1t truy\u1ec7n": "- Ch\u01b0a \u0111i s\u00e2u v\u00e0o ch\u01a1i nhi\u1ec1u ng\u01b0\u1eddi v\u00e0 c\u1ed1t truy\u1ec7n",
    "Tham kh\u1ea3o v\u00e0 l\u1ea5y m\u00f4 h\u00ecnh 3D": "Tham kh\u1ea3o v\u00e0 l\u1ea5y m\u00f4 h\u00ecnh 3D.",
    "Ch\u1ec9nh s\u1eeda m\u00f4 h\u00ecnh 3D v\u00e0 l\u00e0m ho\u1ea1t \u1ea3nh": "Ch\u1ec9nh s\u1eeda m\u00f4 h\u00ecnh 3D v\u00e0 l\u00e0m ho\u1ea1t \u1ea3nh.",
    "Ch\u1ec9nh s\u1eeda h\u00ecnh \u1ea3nh, giao di\u1ec7n": "Ch\u1ec9nh s\u1eeda h\u00ecnh \u1ea3nh v\u00e0 t\u00e0i nguy\u00ean giao di\u1ec7n.",
    "T\u1ed5 ch\u1ee9c d\u1ef1 \u00e1n v\u00e0 nh\u00f3m ch\u1ee9c n\u0103ng": "T\u1ed5 ch\u1ee9c m\u00e3 ngu\u1ed3n v\u00e0 nh\u00f3m ch\u1ee9c n\u0103ng",
    "- Qu\u1ea3n l\u00fd kinh nghi\u1ec7m, l\u00ean c\u1ea5p n\u00e2ng c\u1ea5p v\u00e0 b\u1ea3n \u0111\u1ed3.": "- Qu\u1ea3n l\u00fd kinh nghi\u1ec7m, l\u00ean c\u1ea5p, n\u00e2ng c\u1ea5p v\u00e0 b\u1ea3n \u0111\u1ed3.",
    "Qu\u00e1i, boss v\u00e0 c\u00e1c \u0111\u1ee3t qu\u00e1i": "Qu\u00e1i, boss v\u00e0 c\u00e1c \u0111\u1ee3t t\u1ea5n c\u00f4ng",
    "- C\u00e1c k\u1ebb \u0111\u1ecbch xu\u1ea5t hi\u1ec7n theo t\u1eebng \u0111\u1ee3t qu\u00e1i \u0111\u1ec3 t\u1ea1o \u00e1p l\u1ef1c cho ng\u01b0\u1eddi ch\u01a1i.": "- C\u00e1c lo\u1ea1i qu\u00e1i xu\u1ea5t hi\u1ec7n theo t\u1eebng \u0111\u1ee3t \u0111\u1ec3 t\u1ea1o \u00e1p l\u1ef1c cho ng\u01b0\u1eddi ch\u01a1i.",
    "- Ho\u00e0n th\u00e0nh \u0111\u1ee3t c\u00f3 boss s\u1ebd \u0111\u1ed5i b\u1ea3n \u0111\u1ed3 m\u1edbi": "- Ho\u00e0n th\u00e0nh \u0111\u1ee3t c\u00f3 boss s\u1ebd chuy\u1ec3n sang b\u1ea3n \u0111\u1ed3 m\u1edbi.",
    "giao di\u1ec7n th\u00f4ng tin trong tr\u1eadn, T\u1ea1m d\u1eebng, C\u00e0i \u0110\u1eb7t, B\u1ea3n \u0110\u1ed3": "Giao di\u1ec7n trong tr\u1eadn, t\u1ea1m d\u1eebng v\u00e0 b\u1ea3n \u0111\u1ed3",
    "- giao di\u1ec7n th\u00f4ng tin trong tr\u1eadn hi\u1ec3n th\u1ecb c\u00e1c th\u00f4ng tin ch\u00ednh nh\u01b0 m\u00e1u, kinh nghi\u1ec7m, c\u1ea5p \u0111\u1ed9 v\u00e0 \u0111\u1ee3t qu\u00e1i hi\u1ec7n t\u1ea1i.": "- Giao di\u1ec7n trong tr\u1eadn hi\u1ec3n th\u1ecb m\u00e1u, kinh nghi\u1ec7m, c\u1ea5p \u0111\u1ed9 v\u00e0 \u0111\u1ee3t qu\u00e1i hi\u1ec7n t\u1ea1i.",
    "- B\u1ea3n \u0111\u1ed3 theo giai \u0111o\u1ea1n \u0111\u01b0\u1ee3c thay \u0111\u1ed5i theo t\u1eebng giai \u0111o\u1ea1n \u0111\u1ec3 t\u1ea1o c\u1ea3m gi\u00e1c m\u1edbi trong qu\u00e1 tr\u00ecnh ch\u01a1i.": "- B\u1ea3n \u0111\u1ed3 thay \u0111\u1ed5i theo t\u1eebng giai \u0111o\u1ea1n \u0111\u1ec3 t\u1ea1o c\u1ea3m gi\u00e1c m\u1edbi trong qu\u00e1 tr\u00ecnh ch\u01a1i.",
    "- Ng\u01b0\u1eddi ch\u01a1i t\u01b0\u01a1ng t\u00e1c v\u1edbi khu v\u1ef1c khu v\u1ef1c b\u1eaft \u0111\u1ea7u tr\u1eadn \u0111\u1ec3 m\u1edf b\u1ea3ng h\u01b0\u1edbng d\u1eabn b\u1eaft \u0111\u1ea7u tr\u1eadn.": "- Ng\u01b0\u1eddi ch\u01a1i t\u01b0\u01a1ng t\u00e1c v\u1edbi khu v\u1ef1c b\u1eaft \u0111\u1ea7u tr\u1eadn \u0111\u1ec3 m\u1edf b\u1ea3ng h\u01b0\u1edbng d\u1eabn.",
}


def main():
    prs = Presentation(str(PPTX))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text in REPLACEMENTS:
                        run.text = REPLACEMENTS[run.text]

    prs.save(str(PPTX))
    print(PPTX)


if __name__ == "__main__":
    main()
